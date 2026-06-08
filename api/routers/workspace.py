"""
Router Workspace — Gestion de projets, documents et base de connaissances.

Prefix : /api/workspace
Tables : workspace_projects, workspace_documents, workspace_knowledge
"""

from __future__ import annotations

import json
import logging
import os
import uuid
from datetime import datetime, timezone
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, Query
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from sqlalchemy import text

from api.database import get_session_factory
from api.dependencies import get_db, get_optional_user
from sylea.core.storage.database import DatabaseManager

logger = logging.getLogger("workspace")

router = APIRouter(prefix="/api/workspace", tags=["workspace"])

# ── Export directory ─────────────────────────────────────────────────────────
EXPORT_DIR = Path(__file__).resolve().parent.parent.parent / "data" / "workspace_exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


# ── Schemas ──────────────────────────────────────────────────────────────────

class ProjectIn(BaseModel):
    name: str
    description: str = ""
    category: str = ""


class ProjectOut(BaseModel):
    id: str
    name: str
    description: str
    category: str
    created_at: str
    updated_at: str


class ProjectDetailOut(ProjectOut):
    documents: list[dict] = []


class DocumentIn(BaseModel):
    title: str
    doc_type: str = "note"  # pdf/docx/xlsx/pptx/note
    content_json: str = "{}"
    filepath: str = ""
    filesize: int = 0
    tags: str = ""
    is_template: bool = False


class DocumentOut(BaseModel):
    id: str
    project_id: str | None = None
    title: str
    doc_type: str
    content_json: str
    filepath: str
    filesize: int
    tags: str
    is_template: bool
    created_at: str
    updated_at: str


class KnowledgeIn(BaseModel):
    category: str = "note"  # research/contact/analysis/note
    title: str
    content: str
    source: str = ""


class KnowledgeOut(BaseModel):
    id: str
    category: str
    title: str
    content: str
    source: str
    created_at: str


class ExportIn(BaseModel):
    format: str = "pdf"  # pdf/docx/xlsx/pptx


# ── DB schema init ───────────────────────────────────────────────────────────

def _ensure_workspace_tables(db: DatabaseManager):
    """Create workspace tables if they don't exist."""
    db.conn.execute("""
        CREATE TABLE IF NOT EXISTS workspace_projects (
            id TEXT PRIMARY KEY,
            auth_user_id TEXT,
            name TEXT NOT NULL,
            description TEXT DEFAULT '',
            category TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
    """)
    db.conn.execute("""
        CREATE TABLE IF NOT EXISTS workspace_documents (
            id TEXT PRIMARY KEY,
            project_id TEXT,
            auth_user_id TEXT,
            title TEXT NOT NULL,
            doc_type TEXT DEFAULT 'note',
            content_json TEXT DEFAULT '{}',
            filepath TEXT DEFAULT '',
            filesize INTEGER DEFAULT 0,
            tags TEXT DEFAULT '',
            is_template INTEGER DEFAULT 0,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
    """)
    db.conn.execute("""
        CREATE TABLE IF NOT EXISTS workspace_knowledge (
            id TEXT PRIMARY KEY,
            auth_user_id TEXT,
            category TEXT DEFAULT 'note',
            title TEXT NOT NULL,
            content TEXT DEFAULT '',
            source TEXT DEFAULT '',
            created_at TEXT NOT NULL
        )
    """)
    db.conn.commit()


def _gen_id() -> str:
    return uuid.uuid4().hex[:12]


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


# ── Sanitize helper for PDF export ───────────────────────────────────────────

def _sanitize_text(text: str) -> str:
    """Replace accented characters for fpdf2 compatibility."""
    replacements = {
        "\u00e9": "e", "\u00e8": "e", "\u00ea": "e", "\u00eb": "e",
        "\u00e0": "a", "\u00e2": "a", "\u00e4": "a",
        "\u00f4": "o", "\u00f6": "o",
        "\u00fb": "u", "\u00fc": "u", "\u00f9": "u",
        "\u00ee": "i", "\u00ef": "i",
        "\u00e7": "c",
        "\u00c9": "E", "\u00c8": "E", "\u00ca": "E",
        "\u00c0": "A", "\u00c2": "A",
        "\u00d4": "O", "\u00d6": "O",
        "\u00db": "U", "\u00dc": "U",
        "\u00ce": "I", "\u00cf": "I",
        "\u00c7": "C",
        "\u0153": "oe", "\u0152": "OE",
        "\u00ab": "<<", "\u00bb": ">>",
        "\u2019": "'", "\u2018": "'",
        "\u201c": '"', "\u201d": '"',
        "\u2013": "-", "\u2014": "--",
        "\u2026": "...",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


# ── Export engine ────────────────────────────────────────────────────────────

def _export_to_pdf(content: dict, filename: str) -> Path:
    """Export content_json to PDF using fpdf2."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    title = _sanitize_text(content.get("title", "Document"))
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, title, ln=True, align="C")
    pdf.ln(5)

    sections = content.get("sections", [])
    for section in sections:
        heading = _sanitize_text(section.get("heading", ""))
        if heading:
            pdf.set_font("Helvetica", "B", 12)
            pdf.cell(0, 8, heading, ln=True)
            pdf.ln(2)

        body = _sanitize_text(section.get("body", ""))
        if body:
            pdf.set_font("Helvetica", "", 10)
            pdf.multi_cell(0, 5, body)
            pdf.ln(3)

        bullets = section.get("bullets", [])
        if bullets:
            pdf.set_font("Helvetica", "", 10)
            for bullet in bullets:
                pdf.cell(5)
                pdf.multi_cell(0, 5, "- " + _sanitize_text(bullet))
            pdf.ln(2)

        table = section.get("table")
        if table and isinstance(table, list) and len(table) > 0:
            pdf.set_font("Helvetica", "", 9)
            col_count = len(table[0]) if isinstance(table[0], list) else 1
            col_w = (pdf.w - 20) / max(col_count, 1)
            for row in table:
                if isinstance(row, list):
                    for cell in row:
                        pdf.cell(col_w, 6, _sanitize_text(str(cell)), border=1)
                    pdf.ln()
            pdf.ln(3)

    filepath = EXPORT_DIR / filename
    pdf.output(str(filepath))
    return filepath


def _export_to_docx(content: dict, filename: str) -> Path:
    """Export content_json to DOCX using python-docx."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    title = content.get("title", "Document")
    doc.add_heading(title, level=0)

    for section in content.get("sections", []):
        heading = section.get("heading", "")
        if heading:
            doc.add_heading(heading, level=1)

        body = section.get("body", "")
        if body:
            doc.add_paragraph(body)

        bullets = section.get("bullets", [])
        for bullet in bullets:
            doc.add_paragraph(bullet, style="List Bullet")

        table_data = section.get("table")
        if table_data and isinstance(table_data, list) and len(table_data) > 0:
            cols = len(table_data[0]) if isinstance(table_data[0], list) else 1
            table = doc.add_table(rows=len(table_data), cols=cols)
            table.style = "Table Grid"
            for i, row in enumerate(table_data):
                if isinstance(row, list):
                    for j, cell_val in enumerate(row):
                        table.cell(i, j).text = str(cell_val)

    filepath = EXPORT_DIR / filename
    doc.save(str(filepath))
    return filepath


def _export_to_xlsx(content: dict, filename: str) -> Path:
    """Export content_json to XLSX using openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = content.get("title", "Document")[:31]

    row_idx = 1
    ws.cell(row=row_idx, column=1, value=content.get("title", "Document")).font = Font(bold=True, size=14)
    row_idx += 2

    for section in content.get("sections", []):
        heading = section.get("heading", "")
        if heading:
            ws.cell(row=row_idx, column=1, value=heading).font = Font(bold=True, size=12)
            row_idx += 1

        body = section.get("body", "")
        if body:
            ws.cell(row=row_idx, column=1, value=body)
            row_idx += 1

        bullets = section.get("bullets", [])
        for bullet in bullets:
            ws.cell(row=row_idx, column=1, value="- " + bullet)
            row_idx += 1

        table_data = section.get("table")
        if table_data and isinstance(table_data, list):
            for row in table_data:
                if isinstance(row, list):
                    for j, cell_val in enumerate(row):
                        ws.cell(row=row_idx, column=j + 1, value=str(cell_val))
                    row_idx += 1

        row_idx += 1

    filepath = EXPORT_DIR / filename
    wb.save(str(filepath))
    return filepath


def _export_to_pptx(content: dict, filename: str) -> Path:
    """Export content_json to PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = content.get("title", "Document")

    # One slide per section
    for section in content.get("sections", []):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        heading = section.get("heading", "")
        if slide.shapes.title:
            slide.shapes.title.text = heading

        body_shape = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()

            body_text = section.get("body", "")
            if body_text:
                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                p.text = body_text
                p.font.size = Pt(14)

            bullets = section.get("bullets", [])
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
                p.font.size = Pt(12)

    filepath = EXPORT_DIR / filename
    prs.save(str(filepath))
    return filepath


_EXPORTERS = {
    "pdf": _export_to_pdf,
    "docx": _export_to_docx,
    "xlsx": _export_to_xlsx,
    "pptx": _export_to_pptx,
}


# ── Projects endpoints ───────────────────────────────────────────────────────

@router.get("/projects", response_model=list[ProjectOut])
async def list_projects(
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """List all projects for the current user."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, name, description, category, created_at, updated_at "
                    "FROM workspace_projects WHERE auth_user_id = :uid "
                    "ORDER BY updated_at DESC"
                ),
                {"uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, name, description, category, created_at, updated_at "
                    "FROM workspace_projects ORDER BY updated_at DESC"
                )
            )
        rows = result.mappings().all()
    return [
        ProjectOut(
            id=r["id"], name=r["name"], description=r["description"],
            category=r["category"], created_at=r["created_at"], updated_at=r["updated_at"],
        )
        for r in rows
    ]


@router.post("/projects", response_model=ProjectOut)
async def create_project(
    data: ProjectIn,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Create a new project."""
    _ensure_workspace_tables(db)
    project_id = _gen_id()
    now = _now()
    factory = get_session_factory()
    async with factory() as session:
        try:
            await session.execute(
                text(
                    "INSERT INTO workspace_projects "
                    "(id, auth_user_id, name, description, category, created_at, updated_at) "
                    "VALUES (:id, :uid, :name, :description, :category, :created_at, :updated_at)"
                ),
                {
                    "id": project_id,
                    "uid": user_id or "",
                    "name": data.name,
                    "description": data.description,
                    "category": data.category,
                    "created_at": now,
                    "updated_at": now,
                },
            )
            await session.commit()
        except Exception:
            await session.rollback()
            raise
    return ProjectOut(
        id=project_id, name=data.name, description=data.description,
        category=data.category, created_at=now, updated_at=now,
    )


@router.get("/projects/{project_id}", response_model=ProjectDetailOut)
async def get_project(
    project_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Get project detail with its documents."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, name, description, category, created_at, updated_at "
                    "FROM workspace_projects WHERE id = :pid AND auth_user_id = :uid"
                ),
                {"pid": project_id, "uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, name, description, category, created_at, updated_at "
                    "FROM workspace_projects WHERE id = :pid"
                ),
                {"pid": project_id},
            )
        row = result.mappings().first()
        if not row:
            raise HTTPException(status_code=404, detail="Projet non trouve")

        docs_result = await session.execute(
            text(
                "SELECT id, project_id, auth_user_id, title, doc_type, content_json, filepath, "
                "filesize, tags, is_template, created_at, updated_at "
                "FROM workspace_documents WHERE project_id = :pid ORDER BY updated_at DESC"
            ),
            {"pid": project_id},
        )
        docs = docs_result.mappings().all()

    return ProjectDetailOut(
        id=row["id"], name=row["name"], description=row["description"],
        category=row["category"], created_at=row["created_at"], updated_at=row["updated_at"],
        documents=[
            dict(
                id=d["id"], project_id=d["project_id"], title=d["title"],
                doc_type=d["doc_type"], content_json=d["content_json"],
                filepath=d["filepath"], filesize=d["filesize"], tags=d["tags"],
                is_template=bool(d["is_template"]), created_at=d["created_at"],
                updated_at=d["updated_at"],
            )
            for d in docs
        ],
    )


@router.delete("/projects/{project_id}")
async def delete_project(
    project_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Delete a project and all its documents."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        try:
            if user_id:
                await session.execute(
                    text(
                        "DELETE FROM workspace_documents "
                        "WHERE project_id = :pid AND auth_user_id = :uid"
                    ),
                    {"pid": project_id, "uid": user_id},
                )
                await session.execute(
                    text(
                        "DELETE FROM workspace_projects "
                        "WHERE id = :pid AND auth_user_id = :uid"
                    ),
                    {"pid": project_id, "uid": user_id},
                )
            else:
                await session.execute(
                    text("DELETE FROM workspace_documents WHERE project_id = :pid"),
                    {"pid": project_id},
                )
                await session.execute(
                    text("DELETE FROM workspace_projects WHERE id = :pid"),
                    {"pid": project_id},
                )
            await session.commit()
        except Exception:
            await session.rollback()
            raise
    return {"success": True, "message": "Projet supprime"}


# ── Documents endpoints ──────────────────────────────────────────────────────

@router.post("/projects/{project_id}/documents", response_model=DocumentOut)
async def add_document_to_project(
    project_id: str,
    data: DocumentIn,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Add a document to a project."""
    _ensure_workspace_tables(db)
    doc_id = _gen_id()
    now = _now()
    factory = get_session_factory()
    async with factory() as session:
        try:
            await session.execute(
                text(
                    "INSERT INTO workspace_documents "
                    "(id, project_id, auth_user_id, title, doc_type, content_json, filepath, "
                    "filesize, tags, is_template, created_at, updated_at) "
                    "VALUES (:id, :pid, :uid, :title, :doc_type, :content_json, :filepath, "
                    ":filesize, :tags, :is_template, :created_at, :updated_at)"
                ),
                {
                    "id": doc_id,
                    "pid": project_id,
                    "uid": user_id or "",
                    "title": data.title,
                    "doc_type": data.doc_type,
                    "content_json": data.content_json,
                    "filepath": data.filepath,
                    "filesize": data.filesize,
                    "tags": data.tags,
                    "is_template": 1 if data.is_template else 0,
                    "created_at": now,
                    "updated_at": now,
                },
            )
            # Update project's updated_at
            await session.execute(
                text("UPDATE workspace_projects SET updated_at = :now WHERE id = :pid"),
                {"now": now, "pid": project_id},
            )
            await session.commit()
        except Exception:
            await session.rollback()
            raise
    return DocumentOut(
        id=doc_id, project_id=project_id, title=data.title, doc_type=data.doc_type,
        content_json=data.content_json, filepath=data.filepath, filesize=data.filesize,
        tags=data.tags, is_template=data.is_template, created_at=now, updated_at=now,
    )


@router.get("/documents", response_model=list[DocumentOut])
async def list_all_documents(
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """List all documents across all projects for the current user."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE auth_user_id = :uid AND is_template = 0 "
                    "ORDER BY updated_at DESC"
                ),
                {"uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE is_template = 0 ORDER BY updated_at DESC"
                )
            )
        rows = result.mappings().all()
    return [
        DocumentOut(
            id=r["id"], project_id=r["project_id"], title=r["title"],
            doc_type=r["doc_type"], content_json=r["content_json"],
            filepath=r["filepath"], filesize=r["filesize"], tags=r["tags"],
            is_template=bool(r["is_template"]), created_at=r["created_at"],
            updated_at=r["updated_at"],
        )
        for r in rows
    ]


@router.get("/documents/{doc_id}", response_model=DocumentOut)
async def get_document(
    doc_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Get a single document by ID."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE id = :did AND auth_user_id = :uid"
                ),
                {"did": doc_id, "uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE id = :did"
                ),
                {"did": doc_id},
            )
        row = result.mappings().first()
    if not row:
        raise HTTPException(status_code=404, detail="Document non trouve")
    return DocumentOut(
        id=row["id"], project_id=row["project_id"], title=row["title"],
        doc_type=row["doc_type"], content_json=row["content_json"],
        filepath=row["filepath"], filesize=row["filesize"], tags=row["tags"],
        is_template=bool(row["is_template"]), created_at=row["created_at"],
        updated_at=row["updated_at"],
    )


@router.delete("/documents/{doc_id}")
async def delete_document(
    doc_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Delete a document."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        try:
            if user_id:
                await session.execute(
                    text(
                        "DELETE FROM workspace_documents "
                        "WHERE id = :did AND auth_user_id = :uid"
                    ),
                    {"did": doc_id, "uid": user_id},
                )
            else:
                await session.execute(
                    text("DELETE FROM workspace_documents WHERE id = :did"),
                    {"did": doc_id},
                )
            await session.commit()
        except Exception:
            await session.rollback()
            raise
    return {"success": True, "message": "Document supprime"}


@router.post("/documents/{doc_id}/duplicate", response_model=DocumentOut)
async def duplicate_document(
    doc_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Duplicate a document."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        try:
            if user_id:
                result = await session.execute(
                    text(
                        "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                        "filepath, filesize, tags, is_template, created_at, updated_at "
                        "FROM workspace_documents WHERE id = :did AND auth_user_id = :uid"
                    ),
                    {"did": doc_id, "uid": user_id},
                )
            else:
                result = await session.execute(
                    text(
                        "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                        "filepath, filesize, tags, is_template, created_at, updated_at "
                        "FROM workspace_documents WHERE id = :did"
                    ),
                    {"did": doc_id},
                )
            row = result.mappings().first()
            if not row:
                raise HTTPException(status_code=404, detail="Document non trouve")

            new_id = _gen_id()
            now = _now()
            new_title = row["title"] + " (copie)"
            await session.execute(
                text(
                    "INSERT INTO workspace_documents "
                    "(id, project_id, auth_user_id, title, doc_type, content_json, filepath, "
                    "filesize, tags, is_template, created_at, updated_at) "
                    "VALUES (:id, :pid, :uid, :title, :doc_type, :content_json, :filepath, "
                    ":filesize, :tags, :is_template, :created_at, :updated_at)"
                ),
                {
                    "id": new_id,
                    "pid": row["project_id"],
                    "uid": user_id or row["auth_user_id"],
                    "title": new_title,
                    "doc_type": row["doc_type"],
                    "content_json": row["content_json"],
                    "filepath": "",
                    "filesize": row["filesize"],
                    "tags": row["tags"],
                    "is_template": row["is_template"],
                    "created_at": now,
                    "updated_at": now,
                },
            )
            await session.commit()
        except HTTPException:
            await session.rollback()
            raise
        except Exception:
            await session.rollback()
            raise
    return DocumentOut(
        id=new_id, project_id=row["project_id"], title=new_title,
        doc_type=row["doc_type"], content_json=row["content_json"],
        filepath="", filesize=row["filesize"], tags=row["tags"],
        is_template=bool(row["is_template"]), created_at=now, updated_at=now,
    )


@router.post("/documents/{doc_id}/export")
async def export_document(
    doc_id: str,
    data: ExportIn,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Export a document to a specific format (pdf/docx/xlsx/pptx)."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE id = :did AND auth_user_id = :uid"
                ),
                {"did": doc_id, "uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE id = :did"
                ),
                {"did": doc_id},
            )
        row = result.mappings().first()
    if not row:
        raise HTTPException(status_code=404, detail="Document non trouve")

    fmt = data.format.lower()
    if fmt not in _EXPORTERS:
        raise HTTPException(status_code=400, detail=f"Format non supporte: {fmt}")

    try:
        content = json.loads(row["content_json"])
    except (json.JSONDecodeError, TypeError):
        content = {"title": row["title"], "sections": [{"heading": "", "body": row["content_json"] or ""}]}

    safe_title = "".join(c for c in row["title"] if c.isalnum() or c in " _-").strip().replace(" ", "_")
    filename = f"{safe_title}_{doc_id}.{fmt}"

    exporter = _EXPORTERS[fmt]
    filepath = exporter(content, filename)

    media_types = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }

    return FileResponse(
        path=str(filepath),
        filename=filename,
        media_type=media_types.get(fmt, "application/octet-stream"),
    )


# ── Templates endpoints ─────────────────────────────────────────────────────

@router.get("/templates", response_model=list[DocumentOut])
async def list_templates(
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """List all template documents."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE auth_user_id = :uid AND is_template = 1 "
                    "ORDER BY updated_at DESC"
                ),
                {"uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                    "filepath, filesize, tags, is_template, created_at, updated_at "
                    "FROM workspace_documents WHERE is_template = 1 ORDER BY updated_at DESC"
                )
            )
        rows = result.mappings().all()
    return [
        DocumentOut(
            id=r["id"], project_id=r["project_id"], title=r["title"],
            doc_type=r["doc_type"], content_json=r["content_json"],
            filepath=r["filepath"], filesize=r["filesize"], tags=r["tags"],
            is_template=True, created_at=r["created_at"], updated_at=r["updated_at"],
        )
        for r in rows
    ]


@router.post("/templates/from/{doc_id}", response_model=DocumentOut)
async def save_as_template(
    doc_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Save an existing document as a template."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        try:
            if user_id:
                result = await session.execute(
                    text(
                        "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                        "filepath, filesize, tags, is_template, created_at, updated_at "
                        "FROM workspace_documents WHERE id = :did AND auth_user_id = :uid"
                    ),
                    {"did": doc_id, "uid": user_id},
                )
            else:
                result = await session.execute(
                    text(
                        "SELECT id, project_id, auth_user_id, title, doc_type, content_json, "
                        "filepath, filesize, tags, is_template, created_at, updated_at "
                        "FROM workspace_documents WHERE id = :did"
                    ),
                    {"did": doc_id},
                )
            row = result.mappings().first()
            if not row:
                raise HTTPException(status_code=404, detail="Document non trouve")

            new_id = _gen_id()
            now = _now()
            new_title = row["title"] + " (template)"
            await session.execute(
                text(
                    "INSERT INTO workspace_documents "
                    "(id, project_id, auth_user_id, title, doc_type, content_json, filepath, "
                    "filesize, tags, is_template, created_at, updated_at) "
                    "VALUES (:id, :pid, :uid, :title, :doc_type, :content_json, :filepath, "
                    ":filesize, :tags, 1, :created_at, :updated_at)"
                ),
                {
                    "id": new_id,
                    "pid": row["project_id"],
                    "uid": user_id or row["auth_user_id"],
                    "title": new_title,
                    "doc_type": row["doc_type"],
                    "content_json": row["content_json"],
                    "filepath": "",
                    "filesize": row["filesize"],
                    "tags": row["tags"],
                    "created_at": now,
                    "updated_at": now,
                },
            )
            await session.commit()
        except HTTPException:
            await session.rollback()
            raise
        except Exception:
            await session.rollback()
            raise
    return DocumentOut(
        id=new_id, project_id=row["project_id"], title=new_title,
        doc_type=row["doc_type"], content_json=row["content_json"],
        filepath="", filesize=row["filesize"], tags=row["tags"],
        is_template=True, created_at=now, updated_at=now,
    )


# ── Knowledge base endpoints ────────────────────────────────────────────────

@router.get("/knowledge", response_model=list[KnowledgeOut])
async def list_knowledge(
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """List all knowledge base entries."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, category, title, content, source, created_at "
                    "FROM workspace_knowledge WHERE auth_user_id = :uid "
                    "ORDER BY created_at DESC"
                ),
                {"uid": user_id},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, category, title, content, source, created_at "
                    "FROM workspace_knowledge ORDER BY created_at DESC"
                )
            )
        rows = result.mappings().all()
    return [
        KnowledgeOut(
            id=r["id"], category=r["category"], title=r["title"],
            content=r["content"], source=r["source"], created_at=r["created_at"],
        )
        for r in rows
    ]


@router.post("/knowledge", response_model=KnowledgeOut)
async def add_knowledge(
    data: KnowledgeIn,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Add a knowledge base entry."""
    _ensure_workspace_tables(db)
    entry_id = _gen_id()
    now = _now()
    factory = get_session_factory()
    async with factory() as session:
        try:
            await session.execute(
                text(
                    "INSERT INTO workspace_knowledge "
                    "(id, auth_user_id, category, title, content, source, created_at) "
                    "VALUES (:id, :uid, :category, :title, :content, :source, :now)"
                ),
                {
                    "id": entry_id,
                    "uid": user_id or "",
                    "category": data.category,
                    "title": data.title,
                    "content": data.content,
                    "source": data.source,
                    "now": now,
                },
            )
            await session.commit()
        except Exception:
            await session.rollback()
            raise
    return KnowledgeOut(
        id=entry_id, category=data.category, title=data.title,
        content=data.content, source=data.source, created_at=now,
    )


@router.delete("/knowledge/{entry_id}")
async def delete_knowledge(
    entry_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Delete a knowledge base entry."""
    _ensure_workspace_tables(db)
    factory = get_session_factory()
    async with factory() as session:
        try:
            if user_id:
                await session.execute(
                    text(
                        "DELETE FROM workspace_knowledge "
                        "WHERE id = :eid AND auth_user_id = :uid"
                    ),
                    {"eid": entry_id, "uid": user_id},
                )
            else:
                await session.execute(
                    text("DELETE FROM workspace_knowledge WHERE id = :eid"),
                    {"eid": entry_id},
                )
            await session.commit()
        except Exception:
            await session.rollback()
            raise
    return {"success": True, "message": "Entree supprimee"}


@router.get("/knowledge/search", response_model=list[KnowledgeOut])
async def search_knowledge(
    q: str = Query(..., min_length=1),
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Search knowledge base entries by title and content (LIKE search)."""
    _ensure_workspace_tables(db)
    pattern = f"%{q}%"
    factory = get_session_factory()
    async with factory() as session:
        if user_id:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, category, title, content, source, created_at "
                    "FROM workspace_knowledge "
                    "WHERE auth_user_id = :uid "
                    "AND (title LIKE :pattern OR content LIKE :pattern OR category LIKE :pattern) "
                    "ORDER BY created_at DESC"
                ),
                {"uid": user_id, "pattern": pattern},
            )
        else:
            result = await session.execute(
                text(
                    "SELECT id, auth_user_id, category, title, content, source, created_at "
                    "FROM workspace_knowledge "
                    "WHERE title LIKE :pattern OR content LIKE :pattern OR category LIKE :pattern "
                    "ORDER BY created_at DESC"
                ),
                {"pattern": pattern},
            )
        rows = result.mappings().all()
    return [
        KnowledgeOut(
            id=r["id"], category=r["category"], title=r["title"],
            content=r["content"], source=r["source"], created_at=r["created_at"],
        )
        for r in rows
    ]


# ── Sharing endpoints ───────────────────────────────────────────────────────

def _ensure_shares_table(db: DatabaseManager):
    """Create workspace_shares table if it doesn't exist."""
    db.conn.execute("""
        CREATE TABLE IF NOT EXISTS workspace_shares (
            id TEXT PRIMARY KEY,
            project_id TEXT NOT NULL,
            owner_user_id TEXT NOT NULL,
            shared_with_user_id TEXT NOT NULL,
            permission TEXT DEFAULT 'read',
            created_at TEXT NOT NULL,
            UNIQUE(project_id, shared_with_user_id)
        )
    """)
    db.conn.commit()


class ShareIn(BaseModel):
    shared_with_user_id: str
    permission: str = "read"  # read / write


class ShareOut(BaseModel):
    id: str
    project_id: str
    owner_user_id: str
    shared_with_user_id: str
    permission: str
    created_at: str


class SharedProjectOut(BaseModel):
    id: str
    name: str
    description: str
    category: str
    owner_user_id: str
    permission: str
    shared_at: str


@router.post("/projects/{project_id}/share", response_model=ShareOut)
async def share_project(
    project_id: str,
    data: ShareIn,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Share a project with another user."""
    _ensure_workspace_tables(db)
    _ensure_shares_table(db)

    if not user_id:
        raise HTTPException(status_code=401, detail="Authentification requise pour partager un projet")

    if data.shared_with_user_id == user_id:
        raise HTTPException(status_code=400, detail="Impossible de partager un projet avec soi-meme")

    if data.permission not in ("read", "write"):
        raise HTTPException(status_code=400, detail="Permission invalide (read ou write)")

    share_id = _gen_id()
    now = _now()
    factory = get_session_factory()
    async with factory() as session:
        # Verify the project exists and belongs to the current user
        proj_result = await session.execute(
            text(
                "SELECT 1 FROM workspace_projects WHERE id = :pid AND auth_user_id = :uid"
            ),
            {"pid": project_id, "uid": user_id},
        )
        if not proj_result.first():
            raise HTTPException(status_code=404, detail="Projet non trouve ou acces refuse")

        try:
            await session.execute(
                text(
                    "INSERT INTO workspace_shares "
                    "(id, project_id, owner_user_id, shared_with_user_id, permission, created_at) "
                    "VALUES (:id, :pid, :owner, :sid, :perm, :now)"
                ),
                {
                    "id": share_id,
                    "pid": project_id,
                    "owner": user_id,
                    "sid": data.shared_with_user_id,
                    "perm": data.permission,
                    "now": now,
                },
            )
            await session.commit()
        except Exception:
            await session.rollback()
            # UNIQUE constraint -- already shared
            raise HTTPException(status_code=409, detail="Projet deja partage avec cet utilisateur")

    return ShareOut(
        id=share_id,
        project_id=project_id,
        owner_user_id=user_id,
        shared_with_user_id=data.shared_with_user_id,
        permission=data.permission,
        created_at=now,
    )


@router.get("/shared", response_model=list[SharedProjectOut])
async def get_shared_projects(
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """List projects shared with the current user."""
    _ensure_workspace_tables(db)
    _ensure_shares_table(db)

    if not user_id:
        return []

    factory = get_session_factory()
    async with factory() as session:
        result = await session.execute(
            text(
                "SELECT s.id AS share_id, s.project_id, s.owner_user_id, s.permission, "
                "s.created_at, p.name, p.description, p.category "
                "FROM workspace_shares s "
                "JOIN workspace_projects p ON p.id = s.project_id "
                "WHERE s.shared_with_user_id = :uid "
                "ORDER BY s.created_at DESC"
            ),
            {"uid": user_id},
        )
        rows = result.mappings().all()

    return [
        SharedProjectOut(
            id=r["project_id"],
            name=r["name"],
            description=r["description"],
            category=r["category"],
            owner_user_id=r["owner_user_id"],
            permission=r["permission"],
            shared_at=r["created_at"],
        )
        for r in rows
    ]


@router.get("/projects/{project_id}/shares", response_model=list[ShareOut])
async def list_project_shares(
    project_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """List all shares for a project (owner only)."""
    _ensure_workspace_tables(db)
    _ensure_shares_table(db)

    if not user_id:
        raise HTTPException(status_code=401, detail="Authentification requise")

    factory = get_session_factory()
    async with factory() as session:
        # Verify ownership
        proj_result = await session.execute(
            text(
                "SELECT 1 FROM workspace_projects WHERE id = :pid AND auth_user_id = :uid"
            ),
            {"pid": project_id, "uid": user_id},
        )
        if not proj_result.first():
            raise HTTPException(status_code=404, detail="Projet non trouve ou acces refuse")

        result = await session.execute(
            text(
                "SELECT id, project_id, owner_user_id, shared_with_user_id, permission, created_at "
                "FROM workspace_shares WHERE project_id = :pid ORDER BY created_at DESC"
            ),
            {"pid": project_id},
        )
        rows = result.mappings().all()

    return [
        ShareOut(
            id=r["id"],
            project_id=r["project_id"],
            owner_user_id=r["owner_user_id"],
            shared_with_user_id=r["shared_with_user_id"],
            permission=r["permission"],
            created_at=r["created_at"],
        )
        for r in rows
    ]


@router.delete("/projects/{project_id}/share/{share_id}")
async def unshare_project(
    project_id: str,
    share_id: str,
    db: DatabaseManager = Depends(get_db),
    user_id: str | None = Depends(get_optional_user),
):
    """Remove a share (owner only)."""
    _ensure_workspace_tables(db)
    _ensure_shares_table(db)

    if not user_id:
        raise HTTPException(status_code=401, detail="Authentification requise")

    factory = get_session_factory()
    async with factory() as session:
        try:
            # Verify ownership
            proj_result = await session.execute(
                text(
                    "SELECT 1 FROM workspace_projects WHERE id = :pid AND auth_user_id = :uid"
                ),
                {"pid": project_id, "uid": user_id},
            )
            if not proj_result.first():
                raise HTTPException(status_code=404, detail="Projet non trouve ou acces refuse")

            del_result = await session.execute(
                text(
                    "DELETE FROM workspace_shares WHERE id = :sid AND project_id = :pid"
                ),
                {"sid": share_id, "pid": project_id},
            )
            await session.commit()

            if del_result.rowcount == 0:
                raise HTTPException(status_code=404, detail="Partage non trouve")
        except HTTPException:
            await session.rollback()
            raise
        except Exception:
            await session.rollback()
            raise

    return {"success": True, "message": "Partage supprime"}
